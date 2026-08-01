from __future__ import annotations

from pathlib import Path


def extract_text(path: Path) -> str:
    suffix = path.suffix.lower()
    try:
        if suffix in {".txt", ".md", ".json", ".csv", ".log", ".py", ".js", ".ts", ".html", ".css"}:
            return path.read_text(encoding="utf-8", errors="replace")
        if suffix == ".docx":
            from docx import Document

            doc = Document(str(path))
            blocks = [paragraph.text.strip() for paragraph in doc.paragraphs if paragraph.text.strip()]
            for table in doc.tables:
                for row in table.rows:
                    values = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                    if any(values):
                        blocks.append(" | ".join(values))
            return "\n\n".join(blocks)
        if suffix == ".pptx":
            from pptx import Presentation

            presentation = Presentation(str(path))
            blocks: list[str] = []
            for index, slide in enumerate(presentation.slides, start=1):
                slide_blocks = [f"DIAPOSITIVA {index}"]
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False):
                        text = "\n".join(paragraph.text.strip() for paragraph in shape.text_frame.paragraphs if paragraph.text.strip())
                        if text:
                            slide_blocks.append(text)
                    if getattr(shape, "has_table", False):
                        for row in shape.table.rows:
                            values = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                            if any(values):
                                slide_blocks.append(" | ".join(values))
                try:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        slide_blocks.append("NOTAS DEL EXPOSITOR\n" + notes)
                except (AttributeError, ValueError):
                    pass
                blocks.append("\n".join(slide_blocks))
            return "\n\n".join(blocks)
        if suffix == ".pdf":
            from pypdf import PdfReader

            reader = PdfReader(str(path))
            return "\n\n".join((page.extract_text() or "") for page in reader.pages)
    except Exception as exc:
        return f"[No se pudo extraer el texto: {exc}]"
    return ""
