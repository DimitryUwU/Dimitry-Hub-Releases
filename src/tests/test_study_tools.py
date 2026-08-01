from __future__ import annotations

import tempfile
import unittest
from pathlib import Path
from unittest.mock import patch

from fastapi import HTTPException


class StudyToolsTests(unittest.TestCase):
    def test_powerpoint_text_is_extracted_by_slide(self):
        from pptx import Presentation
        from app.extractors import extract_text

        with tempfile.TemporaryDirectory() as tmp:
            path = Path(tmp) / "clase.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = "Derecho administrativo"
            slide.placeholders[1].text = "El acto administrativo produce efectos jurídicos.\nDebe estar motivado."
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = "Elementos"
            slide.placeholders[1].text = "Competencia\nObjeto\nFinalidad pública"
            presentation.save(path)

            text = extract_text(path)
            self.assertIn("DIAPOSITIVA 1", text)
            self.assertIn("Derecho administrativo", text)
            self.assertIn("DIAPOSITIVA 2", text)
            self.assertIn("Finalidad pública", text)

    def test_word_tables_are_included_in_extraction(self):
        from docx import Document
        from app.extractors import extract_text

        with tempfile.TemporaryDirectory() as tmp:
            path = Path(tmp) / "monografia.docx"
            document = Document()
            document.add_heading("Tema de prueba", level=1)
            document.add_paragraph("La motivación explica las razones de una decisión.")
            table = document.add_table(rows=2, cols=2)
            table.cell(0, 0).text = "Concepto"
            table.cell(0, 1).text = "Definición"
            table.cell(1, 0).text = "Competencia"
            table.cell(1, 1).text = "Atribución para actuar"
            document.save(path)

            text = extract_text(path)
            self.assertIn("Tema de prueba", text)
            self.assertIn("Concepto | Definición", text)
            self.assertIn("Competencia | Atribución para actuar", text)

    def test_ficha_is_hierarchical_and_source_faithful(self):
        from app.study_tools import generate_study_material

        source = """DIAPOSITIVA 1
Acto administrativo
El acto administrativo produce efectos jurídicos individuales. Debe estar motivado.

DIAPOSITIVA 2
Requisitos
La competencia es la atribución legal para actuar. El objeto debe ser lícito y posible.
"""
        result = generate_study_material("ficha", source)
        response = result["response"]
        self.assertIn("# Ficha de estudio", response)
        self.assertIn("## Contenido organizado por temas", response)
        self.assertIn("## Preguntas de repaso", response)
        self.assertIn("Acto administrativo", response)
        self.assertIn("efectos jurídicos individuales", response)
        self.assertIn("no agrega datos", response)
        self.assertEqual(2, result["section_count"])

    def test_all_study_formats_work_without_ai(self):
        from app.study_tools import generate_study_material

        source = "TEMA PRINCIPAL\nLa motivación expresa razones verificables. La competencia permite actuar dentro de una atribución. El objeto debe ser lícito."
        for kind in ("ficha", "simple", "exam", "speech", "cards", "outline"):
            with self.subTest(kind=kind):
                result = generate_study_material(kind, source)
                self.assertGreater(len(result["response"]), 80)
                self.assertEqual(kind, result["kind"])

    def test_study_endpoint_uses_local_fallback_when_ai_is_unavailable(self):
        from app.main import StudyRequest, study_generate

        payload = StudyRequest(
            kind="ficha",
            content="TEMA\nLa motivación expresa las razones de la decisión. La competencia permite actuar dentro de una atribución definida.",
            use_web=False,
        )
        with patch("app.main._ai_call", side_effect=HTTPException(503, "sin proveedor configurado")):
            result = study_generate(payload)
        self.assertEqual("local", result["provider"])
        self.assertIn("# Ficha de estudio", result["response"])
        self.assertIn("sin proveedor configurado", result["warning"])


if __name__ == "__main__":
    unittest.main()
